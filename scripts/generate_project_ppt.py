#!/usr/bin/env python3
"""Generate a 10-slide student project presentation."""

from __future__ import annotations

import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
ASSETS = DOCS / "report_assets"
METRICS_PATH = ROOT / "website" / "models" / "vit_ensemble_metrics.json"
OUTPUT_PPTX = DOCS / "Skin_Disease_Project_Presentation.pptx"


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle + "\n\nPrepared by: Student Group"


def _style_title(slide, size_pt: int = 34) -> None:
    if not slide.shapes.title:
        return
    p = slide.shapes.title.text_frame.paragraphs[0]
    p.font.size = Pt(size_pt)


def _set_bullet_style(paragraph, size_pt: int = 21) -> None:
    paragraph.font.size = Pt(size_pt)
    paragraph.line_spacing = 1.4
    paragraph.space_after = Pt(12)


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    _style_title(slide, 31)

    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.35), Inches(12.1), Inches(5.8))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.05)

    for idx, line in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        _set_bullet_style(p, 21)


def _fit_image_in_box(image_path: Path, box_w: float, box_h: float) -> tuple[float, float]:
    with Image.open(image_path) as img:
        img_w, img_h = img.size

    img_ratio = img_w / img_h
    box_ratio = box_w / box_h
    if img_ratio >= box_ratio:
        w = box_w
        h = box_w / img_ratio
    else:
        h = box_h
        w = box_h * img_ratio
    return w, h


def add_image_slide(
    prs: Presentation,
    title: str,
    image_path: Path,
    caption: str | None = None,
    insights: list[str] | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    _style_title(slide, 31)

    # Left: image panel with aspect-ratio-safe fit (no stretching).
    panel_left = 0.55
    panel_top = 1.35
    panel_w = 8.7
    panel_h = 5.55
    pic_w, pic_h = _fit_image_in_box(image_path, panel_w, panel_h)
    pic_left = panel_left + (panel_w - pic_w) / 2
    pic_top = panel_top + (panel_h - pic_h) / 2

    slide.shapes.add_picture(
        str(image_path),
        Inches(pic_left),
        Inches(pic_top),
        width=Inches(pic_w),
        height=Inches(pic_h),
    )

    # Right: insights panel to avoid empty area and improve readability.
    note_box = slide.shapes.add_textbox(Inches(9.45), Inches(1.35), Inches(3.2), Inches(5.55))
    ntf = note_box.text_frame
    ntf.clear()
    ntf.word_wrap = True
    p0 = ntf.paragraphs[0]
    p0.text = "Key Insights"
    p0.font.size = Pt(18)
    p0.font.bold = True
    p0.space_after = Pt(10)

    for line in insights or []:
        p = ntf.add_paragraph()
        p.text = line
        p.level = 0
        _set_bullet_style(p, 15)

    if caption:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(6.95), Inches(12.1), Inches(0.4))
        tf = box.text_frame
        tf.text = caption
        tf.paragraphs[0].font.size = Pt(13)


def add_metrics_table_slide(prs: Presentation, metrics: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Final Metrics Table"
    _style_title(slide, 31)

    rows = 10
    cols = 4
    table = slide.shapes.add_table(rows, cols, Inches(0.55), Inches(1.3), Inches(12.2), Inches(5.55)).table

    table.columns[0].width = Inches(3.15)
    table.columns[1].width = Inches(2.7)
    table.columns[2].width = Inches(2.7)
    table.columns[3].width = Inches(3.65)

    headers = ["Class", "Precision", "Recall", "Support"]
    for col, h in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(15)

    order = metrics["class_order"]
    per = metrics["per_class"]
    for i, cls in enumerate(order, start=1):
        table.cell(i, 0).text = cls
        table.cell(i, 1).text = f"{per[cls]['precision']:.4f}"
        table.cell(i, 2).text = f"{per[cls]['recall']:.4f}"
        table.cell(i, 3).text = str(per[cls]["support"])
        for c in range(cols):
            table.cell(i, c).text_frame.paragraphs[0].font.size = Pt(14)

    note = slide.shapes.add_textbox(Inches(0.55), Inches(6.9), Inches(12.2), Inches(0.45))
    tf = note.text_frame
    tf.text = (
        f"Overall: Accuracy={metrics['accuracy']:.4f}, Macro Precision={metrics['macro_precision']:.4f}, "
        f"Macro Recall={metrics['macro_recall']:.4f}"
    )
    tf.paragraphs[0].font.size = Pt(13)


def add_roadmap_slide(prs: Presentation) -> None:
    bullets = [
        "Build strict non-overlap benchmark split for final academic evaluation",
        "Train higher-capacity ViT with longer schedule and distillation",
        "Add confidence thresholding and abstention for uncertain cases",
        "Run external validation on a separate dataset",
        "Improve Eczema and Tinea recall via targeted data curation",
        "Deploy versioned monitoring dashboard for drift and class-wise degradation",
    ]
    add_bullets_slide(prs, "Future Scope and Roadmap", bullets)


def main() -> int:
    metrics = json.loads(METRICS_PATH.read_text(encoding="utf-8"))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Skin Disease Classifier",
        "Student Group Project Presentation | Vision Transformer Ensemble",
    )

    add_bullets_slide(
        prs,
        "Problem Statement and Objectives",
        [
            "Classify 8 skin disease categories from images",
            "Develop CPU-friendly ViT-based training and inference pipeline",
            "Increase accuracy using model ensembling and robust evaluation",
            "Integrate best model in web application runtime",
            "Maintain reproducible experiments with saved metrics and confusion outputs",
            "Provide student-ready report and presentation deliverables",
        ],
    )

    add_bullets_slide(
        prs,
        "Dataset and Experimental Setup",
        [
            "Evaluation split: hf_hybrid_150 test set",
            "8 classes, 23 images per class, 184 total test samples",
            "CPU-only training environment",
            "Leakage-safe subset verification used for reliability",
            "Final scoring tracked with accuracy, macro precision, and macro recall",
            "Class-level diagnostics captured using confusion matrix analysis",
        ],
    )

    add_image_slide(
        prs,
        "KPI Summary",
        ASSETS / "kpi_summary.png",
        "Final performance of deployed ensemble",
        [
            "Accuracy reached 85.33% on the test split.",
            "Macro precision and recall both remained above 85%.",
            "Improvement was obtained without GPU training.",
        ],
    )
    add_metrics_table_slide(prs, metrics)
    add_image_slide(
        prs,
        "Class-wise Precision and Recall",
        ASSETS / "precision_recall_by_class.png",
        insights=[
            "Acne and BCC show strong recall.",
            "Eczema recall is the major challenge class.",
            "Per-class balance improved over single-model baseline.",
        ],
    )
    add_image_slide(
        prs,
        "Confusion Matrix",
        ASSETS / "confusion_matrix_heatmap.png",
        insights=[
            "Most errors occur in clinically similar categories.",
            "AK occasionally overlaps with BCC.",
            "Eczema and Psoriasis still need more separation.",
        ],
    )
    add_image_slide(
        prs,
        "Prediction Distribution",
        ASSETS / "true_vs_predicted_counts.png",
        insights=[
            "Predicted volumes are close to true counts.",
            "Some over-prediction appears in Acne and Psoriasis.",
            "Under-prediction in Eczema aligns with recall gap.",
        ],
    )
    add_image_slide(
        prs,
        "Ensemble Composition",
        ASSETS / "ensemble_weights_pie.png",
        insights=[
            "Two strongest members carry 90% total weight.",
            "Small-weight members improve calibration stability.",
            "Weighted voting gave better robustness than single model.",
        ],
    )
    add_roadmap_slide(prs)

    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PPTX))
    print(f"saved: {OUTPUT_PPTX}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
